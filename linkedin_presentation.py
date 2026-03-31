# ============================================================
#  LinkedIn Carousel Presentation Generator
#  Auteur : Malek Bentarbout — Quality Engineer
#  Description : Génère des présentations .pptx professionnelles
#                pour LinkedIn avec palette cream/navy/yellow
# ============================================================

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────
# 🎨 PALETTE DE COULEURS
# ─────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x3E)   # Bleu marine foncé
CREAM      = RGBColor(0xF5, 0xF0, 0xE8)   # Crème clair
YELLOW     = RGBColor(0xF5, 0xC5, 0x18)   # Jaune accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)   # Blanc pur
GRAY       = RGBColor(0x88, 0x88, 0x99)   # Gris texte secondaire

# ─────────────────────────────────────────
# 📐 DIMENSIONS (16:9)
# ─────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─────────────────────────────────────────
# 🔧 FONCTIONS UTILITAIRES
# ─────────────────────────────────────────

def set_background(slide, color):
    """Applique une couleur de fond à un slide."""
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                  font_size=24, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Ajoute une zone de texte stylée sur un slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_yellow_bar(slide, top=Inches(0.15), height=Inches(0.08)):
    """Ajoute une barre jaune décorative en haut du slide."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), top, SLIDE_W, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = YELLOW
    bar.line.fill.background()


def add_branding(slide, text="malekbentarbout"):
    """Ajoute le branding en bas à droite."""
    add_text_box(
        slide, f"© {text}",
        left=Inches(10.5), top=Inches(7.0),
        width=Inches(2.8), height=Inches(0.4),
        font_size=10, color=GRAY, align=PP_ALIGN.RIGHT
    )


# ─────────────────────────────────────────
# 📊 SLIDES
# ─────────────────────────────────────────

def create_cover_slide(prs, title, subtitle, topic_tag):
    """Slide de couverture — fond navy, titre large."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, NAVY)
    add_yellow_bar(slide)

    # Tag sujet (ex: "#FMEA | Quality Engineering")
    add_text_box(slide, topic_tag,
                 left=Inches(0.6), top=Inches(1.2),
                 width=Inches(11), height=Inches(0.5),
                 font_size=14, color=YELLOW, bold=False)

    # Titre principal
    add_text_box(slide, title,
                 left=Inches(0.6), top=Inches(2.0),
                 width=Inches(11), height=Inches(2.0),
                 font_size=44, bold=True, color=WHITE)

    # Sous-titre
    add_text_box(slide, subtitle,
                 left=Inches(0.6), top=Inches(4.2),
                 width=Inches(10), height=Inches(0.8),
                 font_size=18, color=CREAM)

    add_branding(slide)
    return slide


def create_content_slide(prs, slide_number, heading, points):
    """Slide contenu — fond crème, bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, CREAM)
    add_yellow_bar(slide)

    # Numéro du slide
    add_text_box(slide, f"{slide_number:02d}",
                 left=Inches(0.4), top=Inches(0.3),
                 width=Inches(1), height=Inches(0.8),
                 font_size=28, bold=True, color=NAVY)

    # Titre du slide
    add_text_box(slide, heading,
                 left=Inches(1.5), top=Inches(0.35),
                 width=Inches(10), height=Inches(0.9),
                 font_size=26, bold=True, color=NAVY)

    # Ligne de séparation (forme fine)
    line = slide.shapes.add_shape(1,
        Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()

    # Bullet points
    y_pos = Inches(1.7)
    for point in points:
        add_text_box(slide, f"▸  {point}",
                     left=Inches(0.6), top=y_pos,
                     width=Inches(12), height=Inches(0.7),
                     font_size=18, color=NAVY)
        y_pos += Inches(0.85)

    add_branding(slide)
    return slide


def create_closing_slide(prs, cta_text, hashtags):
    """Slide de clôture — fond navy, call to action."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, NAVY)
    add_yellow_bar(slide)

    add_text_box(slide, cta_text,
                 left=Inches(1), top=Inches(2.0),
                 width=Inches(11), height=Inches(2.5),
                 font_size=32, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, hashtags,
                 left=Inches(1), top=Inches(4.8),
                 width=Inches(11), height=Inches(0.8),
                 font_size=14, color=YELLOW,
                 align=PP_ALIGN.CENTER)

    add_branding(slide)
    return slide


# ─────────────────────────────────────────
# 🚀 GÉNÉRATION PRINCIPALE
# ─────────────────────────────────────────

def generate_presentation(filename="linkedin_carousel.pptx"):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── SLIDE 1 : Cover ──────────────────
    create_cover_slide(
        prs,
        title="Les 5 étapes clés du processus FMEA",
        subtitle="Comment anticiper les défaillances avant qu'elles arrivent",
        topic_tag="#FMEA  |  Quality Engineering  |  IATF 16949"
    )

    # ── SLIDES 2–5 : Contenu ─────────────
    slides_content = [
        {
            "heading": "Qu'est-ce que le FMEA ?",
            "points": [
                "FMEA = Failure Mode and Effects Analysis",
                "Outil préventif de management des risques",
                "Obligatoire dans l'industrie automobile (IATF 16949)",
                "3 types : Design FMEA, Process FMEA, System FMEA"
            ]
        },
        {
            "heading": "Étape 1 — Définir le périmètre",
            "points": [
                "Identifier le produit ou le processus à analyser",
                "Constituer une équipe pluridisciplinaire",
                "Définir les frontières et les interfaces",
                "Collecter les données historiques et retours clients"
            ]
        },
        {
            "heading": "Étape 2 — Identifier les modes de défaillance",
            "points": [
                "Lister toutes les fonctions du système",
                "Pour chaque fonction : quels défauts possibles ?",
                "Utiliser les données terrain et l'expérience équipe",
                "Ne pas filtrer à ce stade — exhaustivité d'abord"
            ]
        },
        {
            "heading": "Étape 3 — Évaluer Gravité / Occurrence / Détection",
            "points": [
                "Gravité (S) : impact sur le client — de 1 à 10",
                "Occurrence (O) : probabilité d'apparition — de 1 à 10",
                "Détection (D) : capacité à détecter — de 1 à 10",
                "RPN = S × O × D → priorité des actions"
            ]
        },
    ]

    for i, content in enumerate(slides_content, start=2):
        create_content_slide(prs, i, content["heading"], content["points"])

    # ── SLIDE 6 : Closing ─────────────────
    create_closing_slide(
        prs,
        cta_text="Vous utilisez le FMEA dans votre équipe ?\nPartagez votre expérience en commentaire 👇",
        hashtags="#QualityEngineering  #FMEA  #IATF16949  #Automotive  #LinkedInCarousel"
    )

    # ── Sauvegarde ────────────────────────
    prs.save(filename)
    print(f"✅ Présentation générée : {filename}")


# ─────────────────────────────────────────
# ▶ LANCEMENT
# ─────────────────────────────────────────
if __name__ == "__main__":
    generate_presentation("linkedin_fmea_carousel.pptx")
